# -*- coding: utf-8 -*-
"""
PDF のテキスト層を座標・サイズをできるだけ保って PPTX に配置するスクリプト
使い方:
    python pdf_to_pptx_preserve_text.py [input.pdf]
出力:
    input_preserved.pptx
注意:
    - PDF にテキスト層（selectable text）があることが前提です。
    - 画像・表・縦書き・回転文字はこの簡易版では扱いません。
"""

import sys
from pathlib import Path
from tkinter import Tk, filedialog

import pdfplumber#PDFのテキスト層や文字ごとの位置情報を抽出する
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ---------- 設定（必要に応じ変更） ----------
DEFAULT_FONT_NAME = "Meiryo"      # 環境に合わせて変更
DEFAULT_FONT_COLOR = (0, 0, 0)    # RGB
SLIDE_WIDTH = 10.0             # 出力スライド幅（inch）
SLIDE_HEIGHT = 7.5             # 出力スライド高さ（inch）
LINE_TOP_TOLERANCE_PT = 3.0       # 同一行判定の垂直許容（pt）
CHAR_GAP_FACTOR = 0.5             # 文字間ギャップ判定（size * factor）
MIN_BOX_WIDTH = 0.08           # 最小テキストボックス幅（inch）
MIN_BOX_HEIGHT = 0.08          # 最小テキストボックス高さ（inch）
# ------------------------------------------------

def ask_pdf_path():
    root = Tk(); root.withdraw(); root.update()
    p = filedialog.askopenfilename(title="PDFファイルを選択してください", filetypes=[("PDF Files","*.pdf")])
    root.destroy()
    return p

def group_chars_to_lines(chars, tol=LINE_TOP_TOLERANCE_PT):
    """chars: list of dicts from pdfplumber page.chars
       return: list of (top, [chars_in_line]) sorted by top"""
    if not chars:
        return []
    # sort by top then x0
    chars_sorted = sorted(chars, key=lambda c: (round(c.get("top", 0), 3), c.get("x0", 0)))
    lines = []
    cur_top = chars_sorted[0].get("top", 0)
    cur_line = []
    for c in chars_sorted:
        t = c.get("top", 0)
        if abs(t - cur_top) <= tol:
            cur_line.append(c)
        else:
            lines.append((cur_top, cur_line))
            cur_top = t
            cur_line = [c]
    if cur_line:
        lines.append((cur_top, cur_line))
    # sort chars inside each line by x0
    for i, (t, ln) in enumerate(lines):
        lines[i] = (t, sorted(ln, key=lambda c: c.get("x0", 0)))
    return lines

def build_text_from_line(chars_in_line):
    """chars_in_line: list of char dicts sorted by x0
       returns text string and bounding x0,x1,bottom and avg font size (pt)"""
    if not chars_in_line:
        return "", 0, 0, 0, None
    texts = []
    prev_x1 = None
    sizes = []
    for c in chars_in_line:
        ch = c.get("text", "")
        sizes.append(c.get("size", 0) or 0)
        if prev_x1 is None:
            texts.append(ch)
        else:
            gap = c.get("x0", 0) - prev_x1
            avg_size = (sizes[-1] + (sizes[-2] if len(sizes) >= 2 else sizes[-1])) / 2 if len(sizes) >= 2 else sizes[-1]
            gap_thresh = avg_size * CHAR_GAP_FACTOR if avg_size else 0
            if gap > gap_thresh:
                texts.append(" ")
                texts.append(ch)
            else:
                texts.append(ch)
        prev_x1 = c.get("x1", prev_x1)
    text = "".join(texts)
    x0 = min(c.get("x0", 0) for c in chars_in_line)
    x1 = max(c.get("x1", 0) for c in chars_in_line)
    bottom = max(c.get("bottom", c.get("top", 0) + c.get("size", 0)) for c in chars_in_line)
    sizes_nonzero = [s for s in sizes if s > 0]
    avg_font = sum(sizes_nonzero) / len(sizes_nonzero) if sizes_nonzero else None
    return text, x0, x1, bottom, avg_font

def pdf_page_to_slide(prs, page, slide_w=SLIDE_WIDTH, slide_h=SLIDE_HEIGHT):
    """Convert one pdfplumber page to one slide with positioned textboxes."""
    
    # PDFのページサイズを取得
    pdf_w_pt = page.width
    pdf_h_pt = page.height

    # create slide (use blank layout to avoid placeholder issues)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # group chars into lines
    chars = page.chars
    lines = group_chars_to_lines(chars)

    # scale factors:
    # fraction_x = x_pt / pdf_w_pt  -> left_in = fraction_x * slide_w_in
    # fraction_y = y_pt / pdf_h_pt  -> top_in = fraction_y * slide_h_in
    # For font scaling: PDF font size is in points (pt). To map to PPT point size:
    #   font_pt_ppt = avg_font_pt * (slide_height_in_inches / (pdf_height_in_inches))
    # where pdf_height_in_inches = pdf_h_pt / 72.0
    # So scale_font = slide_h_in / (pdf_h_pt / 72.0) = slide_h_in * 72.0 / pdf_h_pt
    font_scale = slide_h * 72.0 / pdf_h_pt if pdf_h_pt else 1.0

    for top_pt, chars_in_line in lines:
        text, x0, x1, bottom, avg_font = build_text_from_line(chars_in_line)
        if not text.strip():
            continue

        # compute position and size in inches
        left_in = (x0 / pdf_w_pt) * slide_w if pdf_w_pt else 0
        top_in = (top_pt / pdf_h_pt) * slide_h if pdf_h_pt else 0
        width_in = ((x1 - x0) / pdf_w_pt) * slide_w if pdf_w_pt else MIN_BOX_WIDTH
        height_in = ((bottom - top_pt) / pdf_h_pt) * slide_h if pdf_h_pt else MIN_BOX_HEIGHT

        # enforce minimum sizes
        if width_in < MIN_BOX_WIDTH:
            width_in = MIN_BOX_WIDTH
        if height_in < MIN_BOX_HEIGHT:
            height_in = MIN_BOX_HEIGHT

        # add textbox
        left = Inches(left_in)
        top = Inches(top_in)
        width = Inches(width_in)
        height = Inches(height_in)
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text

        # compute font size in points for ppt
        if avg_font:
            font_pt = avg_font * font_scale
            # clamp to reasonable range
            if font_pt < 4:
                font_pt = 4
            if font_pt > 200:
                font_pt = 200
        else:
            font_pt = 12

        # apply font to runs
        for run in p.runs:
            try:
                run.font.size = Pt(font_pt)
                run.font.name = DEFAULT_FONT_NAME
                run.font.color.rgb = RGBColor(*DEFAULT_FONT_COLOR)
                # set eastAsia font name for Japanese if available
                run._element.rPr.rFonts.set(qn('w:eastAsia'), DEFAULT_FONT_NAME)
            except Exception:
                try:
                    run.font.size = Pt(font_pt)
                except Exception:
                    pass

    return slide

def convert_pdf_to_pptx(pdf_path, output_pptx):
    #pptxのプレゼンテーションオブジェクトを作成し、スライドサイズを設定
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    #pdfをスライドごとに処理していく
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            pdf_page_to_slide(prs, page, SLIDE_WIDTH, SLIDE_HEIGHT)

    #pptxファイルとして保存
    prs.save(output_pptx)
    #print("Saved:", output_pptx)

# GUIでPDF を選択させる関数
def ask_pdf_path_gui():
    root = Tk()#ルートウィンドウを作成
    try:
        root.withdraw()#ルートウィンドウを非表示にする
        root.update()  # プラットフォームによっては必要
        selected_pdf = filedialog.askopenfilename(
            title="PDFファイルを選択してください",
            filetypes=[("PDF Files", ("*.pdf","*.PDF"))]
        )
    finally:
        root.destroy()
    return selected_pdf

#入力部分
def main():
    selected_pdf = ask_pdf_path_gui()# GUIでPDF を選択させる

    # PDF が選択されなかった場合は終了
    if not selected_pdf:
        print("PDF が選択されませんでした")
        return

    pdf_path = Path(selected_pdf)# 入力パスを Path オブジェクトに変換
    output_pptx = str(pdf_path.with_name(pdf_path.stem + "_edited.pptx"))#出力するファイル名
    convert_pdf_to_pptx(str(pdf_path), output_pptx)

if __name__ == "__main__":
    main()
